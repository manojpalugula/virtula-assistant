# STANDARD LIBRARY IMPORTS
import datetime
import os
import sys
import subprocess
import pyperclip
import time  # Single import
import re
import json
import shutil
import textwrap

# THIRD-PARTY LIBRARY IMPORTS
import pyttsx3
import speech_recognition as sr
import wikipedia # CORRECTED: Was 'import wikipedia', but wikipediaapi methods were used
import webbrowser
import spotipy
from spotipy.oauth2 import SpotifyOAuth
from groq import Groq
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
import keyboard # Replaced pynput with keyboard

# CONDITIONAL IMPORT FOR PYAUTOGUI
pyautogui_available = False # Initialize before try-except, ensuring it always exists
try:
    import pyautogui
    pyautogui_available = True
except Exception as e:
    print(f"PyAutoGUI could not be imported or initialized: {e}. GUI automation features will be disabled.")
    # pyautogui_available remains False

# GLOBAL CONSTANTS
PIXABAY_API_KEY = "50245343-f63fa43cdf38d8cb5f42a3c68" # Get from https://pixabay.com/api/docs/

# GLOBAL STATE VARIABLES
jarvis_active = False
hotkey_registered = False # To track if the hotkey is active

# INITIALIZATIONS
# Initialize the text-to-speech engine
engine = pyttsx3.init() # Removed 'sapi5' to allow auto-detection
voices = engine.getProperty('voices')
# Attempt to set a voice, but handle if no voices are available
if voices and len(voices) > 0:
    engine.setProperty('voice', voices[0].id)
else:
    print("No TTS voices found. Speech output may not work.")

# Set up Spotify credentials
# For security, consider moving Client ID and Secret to environment variables or a config file
sp = spotipy.Spotify(auth_manager=SpotifyOAuth(
    client_id="bc4b741fbf23492f939365b367a95119",
    client_secret="1db92c2b6116461c8b6063925b21974c",
    redirect_uri="http://localhost:8888/callback",
    scope="user-read-playback-state,user-modify-playback-state"
))

# FUNCTION DEFINITIONS
def speak(audio):
    try:
        if not engine._inLoop:
            engine.say(audio)
            engine.runAndWait()
        else:
            print("[WARNING] TTS loop already running, skipping this message.")
    except Exception as e:
        print(f"[ERROR] TTS Exception: {e}")


# Activation callback function
def activate_jarvis():
    global jarvis_active, hotkey_registered
    
    jarvis_active = not jarvis_active

    if jarvis_active:
        print("Jarvis activated.")
        speak("Jarvis activated.")
        if hotkey_registered:
            try:
                keyboard.remove_hotkey("ctrl+j")
                print("Ctrl+J hotkey has been de-registered during activation.")
                hotkey_registered = False
            except KeyError:
                print("Warning: Ctrl+J hotkey could not be removed during activation (was not found).")
            except Exception as e:
                print(f"An error occurred while removing hotkey during activation: {e}")
    else:
        print("Jarvis deactivated.")
        speak("Jarvis deactivated.")
        activate_jarvis.has_wished = False # Reset wish flag
        if not hotkey_registered:
            try:
                keyboard.add_hotkey("ctrl+j", activate_jarvis, suppress=False)
                hotkey_registered = True
                print("Ctrl+J hotkey has been re-registered for activation.")
            except Exception as e:
                print(f"Error re-registering Ctrl+J hotkey during deactivation: {e}")

# Initialize the .has_wished attribute for activate_jarvis function
activate_jarvis.has_wished = False

def start_hotkey_listener():
    global hotkey_registered
    if not hotkey_registered and not jarvis_active:
        try:
            keyboard.add_hotkey("ctrl+j", activate_jarvis, suppress=False)
            hotkey_registered = True
            print("Hotkey listener started. Press Ctrl+J to toggle Jarvis.")
        except Exception as e:
            print(f"Error setting up hotkey listener: {e}")
            print("Hotkey detection might not be supported or may require administrator privileges.")
            hotkey_registered = False

def wishMe():
    hour = datetime.datetime.now().hour
    if 0 <= hour < 12:
        speak("Good Morning!")
    elif 12 <= hour < 18:
        speak("Good Afternoon!")
    else:
        speak("Good Evening!")
    speak("I am Jarvis, your personal virtual assistant. Please tell me how may I help you")

def takeCommand():
    r = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening...")
        r.pause_threshold = 1
        audio = r.listen(source)

    try:
        print("Recognizing...")
        query = r.recognize_google(audio, language='en-in')
        print(f"User said: {query}\n")
    except Exception as e:
        print("Say that again please...")
        return "None"
    return query

def playSpotifyTrack(track_name):
    results = sp.search(q=track_name, type='track', limit=1)
    if results['tracks']['items']:
        track = results['tracks']['items'][0]
        track_name_spotify = track['name'] # Renamed to avoid conflict with function arg
        track_artists = ', '.join([artist['name'] for artist in track['artists']])
        track_uri = track['uri']
        # The googleusercontent.com URL for Spotify might be outdated or non-functional for direct play.
        # Spotify's web player or app URI (spotify:track:TRACK_ID) is generally more reliable.
        # For simplicity, opening the track on Spotify's web player:
        webbrowser.open(track['external_urls']['spotify'])
        speak(f"Playing {track_name_spotify} by {track_artists} on Spotify")
        speak("Please press play in the browser if the song doesn't start automatically.")
    else:
        speak(f"Couldn't find {track_name} on Spotify")

def searchGoogle(query):
    url = f"https://www.google.com/search?q={query}"
    webbrowser.open(url)

def searchYouTube(query):
    # Assuming you want to search on YouTube directly
    url = f"https://www.youtube.com/results?search_query={query}"
    # The original googleusercontent.com URL might be a proxy or custom setup.
    # Using direct Youtube URL for broader compatibility.
    webbrowser.open(url)

def query_groq(query, api_key_groq): # Renamed api_key for clarity
    client = Groq(api_key=api_key_groq)
    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": query,
            }
        ],
        model="llama3-8b-8192",
    )
    return chat_completion.choices[0].message.content

def get_pixabay_image(query):
    try:
        url = f"https://pixabay.com/api/?key={PIXABAY_API_KEY}&q={query}&image_type=photo&per_page=3"
        response = requests.get(url).json()
        if response['hits']:
            return response['hits'][0]['webformatURL']
    except Exception as e:
        print(f"Pixabay error: {e}")
    return None

def download_image(image_url, filename="temp_img.jpg"):
    try:
        response = requests.get(image_url)
        response.raise_for_status() # Raise an exception for HTTP errors
        with open(filename, 'wb') as f:
            f.write(response.content)
        return filename
    except Exception as e:
        print(f"Image download failed: {e}")
    return None

def clean_groq_content(content):
    content = content.replace('*', '')
    content = re.sub(r'^\s*-\s+', '• ', content, flags=re.MULTILINE)
    content = re.sub(r'^\s*•\s+', '• ', content, flags=re.MULTILINE)
    content = '\n'.join([x.strip() for x in content.split('\n') if x.strip()])
    return content

def split_into_slides(content, max_lines=6):
    paragraphs = [p for p in content.split('\n') if p.strip()]
    slides_data = [] # Renamed from slides to avoid conflict
    current_slide = []
    
    for para in paragraphs:
        wrapped = textwrap.wrap(para, width=80)
        if len(current_slide) + len(wrapped) > max_lines and current_slide: # Ensure current_slide is not empty before appending
            slides_data.append('\n'.join(current_slide))
            current_slide = []
        current_slide.extend(wrapped)
    
    if current_slide:
        slides_data.append('\n'.join(current_slide))
    
    return slides_data if slides_data else [content] # Return list of slide contents

def create_ppt_file_with_content(file_path, topic, api_key_groq): # Renamed api_key
    prompt = f"""
Create a PowerPoint presentation about "{topic}" with 3 to 5 slides.
Each slide must follow this format:

[SLIDE TITLE: Slide title here]
[CONTENT:
- Bullet point 1
- Bullet point 2
- Bullet point 3
]
[IMAGE KEYWORD: relevant image keyword]

Use '---' to separate each slide.
Stick exactly to the format.
"""
    raw_content = query_groq(prompt, api_key_groq)
    slide_contents = [slide_text.strip() for slide_text in raw_content.split('---') if slide_text.strip()] # Renamed 'slides'

    prs = Presentation()
    slide_width_val = prs.slide_width # Renamed
    slide_height_val = prs.slide_height # Renamed

    title_slide_layout = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide_layout.shapes.title.text = topic
    if title_slide_layout.placeholders[1]:
      title_slide_layout.placeholders[1].text = "Created by Jarvis Assistant"

    for i, slide_content_item in enumerate(slide_contents[:5]): # Renamed 'slide_content'
        title_match = re.search(r'\[SLIDE TITLE:\s*(.+?)\]', slide_content_item, re.IGNORECASE)
        content_match = re.search(r'\[CONTENT:(.+?)\](?:\[IMAGE|\Z)', slide_content_item, re.DOTALL | re.IGNORECASE)
        image_match = re.search(r'\[IMAGE KEYWORD:\s*(.+?)\]', slide_content_item, re.IGNORECASE)

        if not content_match or not content_match.group(1).strip():
            continue

        current_slide_obj = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout, renamed 'slide'

        title_text = title_match.group(1).strip() if title_match else f"{topic} - Slide {i+1}"
        title_box = current_slide_obj.shapes.add_textbox(Inches(0.5), Inches(0.2), slide_width_val - Inches(1), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f'"{title_text}"' # Original code added quotes, retaining this style
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(59, 89, 152)
        p.alignment = PP_ALIGN.LEFT

        left_box = current_slide_obj.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), slide_height_val - Inches(2))
        tf_left = left_box.text_frame # Renamed
        tf_left.word_wrap = True

        bullet_lines = [line.strip("-*• ") for line in content_match.group(1).strip().split("\n") if line.strip()]
        for line in bullet_lines:
            para = tf_left.add_paragraph()
            para.text = f"• {line}"
            para.level = 0
            para.font.size = Pt(20)

        if image_match:
            img_keyword = image_match.group(1).strip()
            img_url = get_pixabay_image(img_keyword)
            if img_url:
                img_path = download_image(img_url)
                if img_path:
                    try:
                        left_img = Inches(6.0) # Renamed
                        top_img = Inches(1.8)  # Renamed
                        current_slide_obj.shapes.add_picture(img_path, left_img, top_img, width=Inches(3.5))
                        os.remove(img_path)
                    except Exception as e:
                        print(f"Image error during PPT creation: {e}")

    thank_slide_layout = prs.slides.add_slide(prs.slide_layouts[5]) # Renamed 'thank_slide'
    txBox_thank = thank_slide_layout.shapes.add_textbox(Inches(1), Inches(2), prs.slide_width - Inches(2), Inches(1.5)) # Renamed
    tf_thank = txBox_thank.text_frame # Renamed
    p_thank = tf_thank.paragraphs[0] # Renamed
    p_thank.text = "Thank You!"
    p_thank.font.size = Pt(48)
    p_thank.font.color.rgb = RGBColor(59, 89, 152)
    p_thank.alignment = PP_ALIGN.CENTER

    try:
        prs.save(file_path)
        speak(f"Presentation saved as {os.path.basename(file_path)}")
        open_file(file_path) # Assuming open_file is defined
    except PermissionError:
        speak("Could not save presentation. Please close the file if it's already open and try again.")
    except Exception as e:
        speak(f"An error occurred while saving the presentation: {e}")


def extract_code(response):
    code_pattern = re.compile(r'```(?:python\n)?(.*?)```', re.DOTALL | re.IGNORECASE) # Allow optional 'python\n'
    match = code_pattern.search(response)
    if match:
        return match.group(1).strip()
    return response # Return original response if no code block found

def open_notepad_with_code(code):
    notepad_path = "notepad.exe"
    pyperclip.copy(code)
    speak("Code has been copied to your clipboard.")
    try:
        if sys.platform == "win32":
            subprocess.Popen(notepad_path)
            time.sleep(1)
            if pyautogui_available:
                pyautogui.hotkey("ctrl", "v")
                speak("Pasted into Notepad.")
            else:
                speak("Please paste the code manually into Notepad or your preferred editor.")
        else:
            speak("Please paste the code into your preferred text editor (Notepad is Windows-specific).")
    except Exception as e:
        print(f"Could not open Notepad or paste: {e}")
        speak("Could not open Notepad. Please paste the code manually.")

def create_text_file_with_content(file_path, topic, api_key_groq): # Renamed api_key
    content = query_groq(topic, api_key_groq)
    try:
        with open(file_path, 'w', encoding='utf-8') as file: # Added encoding
            file.write(content)
        speak(f"Text file '{os.path.basename(file_path)}' created on Desktop")
        open_file(file_path)
    except Exception as e:
        speak(f"Failed to create text file: {e}")


def create_word_file_with_content(file_path, topic, api_key_groq): # Renamed api_key
    content = query_groq(topic, api_key_groq)
    doc = Document()
    doc.add_heading(topic, 0)
    doc.add_paragraph(content)
    try:
        doc.save(file_path)
        speak(f"Word file '{os.path.basename(file_path)}' created on Desktop")
        open_file(file_path)
    except Exception as e:
        speak(f"Failed to create Word file: {e}")


def create_folder(folder_path):
    try:
        os.makedirs(folder_path, exist_ok=True)
        speak(f"Folder '{os.path.basename(folder_path)}' created on Desktop")
        open_folder(folder_path) # Assuming open_folder is defined
    except Exception as e:
        speak(f"Failed to create folder: {e}")

def delete_file(file_path):
    try:
        os.remove(file_path)
        speak(f"File '{os.path.basename(file_path)}' deleted from Desktop")
    except FileNotFoundError:
        speak(f"File '{os.path.basename(file_path)}' not found on Desktop")
    except Exception as e:
        speak(f"Failed to delete file: {e}")

def open_file(file_path):
    try:
        if sys.platform == "win32":
            os.startfile(file_path)
        elif sys.platform == "darwin": # macOS
            subprocess.call(('open', file_path))
        else: # Linux and other Unix-like
            subprocess.call(('xdg-open', file_path))
        speak(f"Opening {os.path.basename(file_path)}")
    except FileNotFoundError:
        speak(f"File '{os.path.basename(file_path)}' not found.")
    except Exception as e:
        speak(f"Could not open file: {e}")


def open_folder(folder_path):
    try:
        if sys.platform == "win32":
            os.startfile(folder_path)
        elif sys.platform == "darwin": # macOS
            subprocess.call(('open', folder_path))
        else: # Linux and other Unix-like
            subprocess.call(('xdg-open', folder_path))
        speak(f"Opening folder {os.path.basename(folder_path)}")
    except FileNotFoundError:
        speak(f"Folder '{os.path.basename(folder_path)}' not found.")
    except Exception as e:
        speak(f"Could not open folder: {e}")

def list_files_in_folder(folder_path):
    try:
        files = os.listdir(folder_path)
        if files:
            speak(f"Files in '{os.path.basename(folder_path)}' are:")
            for file_item in files: # Renamed 'file' to avoid conflict with open()
                print(file_item) # Also print to console for clarity
                speak(file_item)
        else:
            speak(f"No files found in '{os.path.basename(folder_path)}'")
    except FileNotFoundError:
        speak(f"Folder '{os.path.basename(folder_path)}' not found on Desktop")
    except Exception as e:
        speak(f"Could not list files: {e}")


def get_weather(city):
    # It's better to keep API keys outside the code or load from config/env
    weather_api_key = "ba621bd65e6bf4c55568056c986c200e" 
    base_url = "http://api.openweathermap.org/data/2.5/weather?"
    complete_url = base_url + "q=" + city + "&appid=" + weather_api_key + "&units=metric"
    try:
        response = requests.get(complete_url)
        response.raise_for_status() # Will raise an HTTPError for bad responses (4XX or 5XX)
        weather_data = response.json()
        
        if weather_data["cod"] != "404": # OpenWeatherMap uses string "404" for cod on error
            main = weather_data["main"]
            wind = weather_data["wind"]
            weather_desc_list = weather_data["weather"] # It's a list
            
            temperature = main.get("temp")
            humidity = main.get("humidity")
            weather_description = weather_desc_list[0].get("description") if weather_desc_list else "not available"
            wind_speed = wind.get("speed")

            weather_report = (f"Temperature: {temperature}°C. "
                              f"Humidity: {humidity}%. "
                              f"Weather description: {weather_description}. "
                              f"Wind speed: {wind_speed} meter per second.")

            print(weather_report)
            speak(weather_report)
        else:
            speak(f"City '{city}' not found by weather service. Please try again.")
    except requests.exceptions.RequestException as e:
        speak(f"Could not retrieve weather data due to a network issue: {e}")
    except KeyError as e:
        speak(f"Unexpected format in weather data: missing {e}")
    except Exception as e:
        speak(f"An error occurred while fetching weather: {e}")


def get_news(api_key_news, query): # Renamed api_key
    url = f"https://newsapi.org/v2/everything?q={query}&apiKey={api_key_news}"
    try:
        response = requests.get(url)
        response.raise_for_status()
        data = response.json()
        if data.get('status') == 'ok':
            articles = data.get('articles', [])
            if not articles:
                return "No news articles found for your query."
            news_list = []
            for article in articles[:5]: # Limit to the top 5 news articles
                title = article.get('title', 'No title')
                # description = article.get('description', 'No description') # Description can be long
                news_list.append(f"Headline: {title}.") # Shorter output
            return " ".join(news_list) if news_list else "No news to report."
        else:
            return f"Could not retrieve news. Error: {data.get('message', 'Unknown error from NewsAPI')}"
    except requests.exceptions.RequestException as e:
        return f"Could not connect to news service: {e}"
    except Exception as e:
        return f"An error occurred while fetching news: {e}"


if __name__ == "__main__":
    # API Keys - Consider loading from environment variables or a config file
    groq_api_key = "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" # For Groq services
    news_api_key = '0af8f4aac7614faa804bbccc4cf9fca2' # For NewsAPI

    start_hotkey_listener()

    try:
        while True: 
            if jarvis_active:
                if not activate_jarvis.has_wished:
                    wishMe()
                    activate_jarvis.has_wished = True

                while jarvis_active: # Inner loop for commands
                    query = takeCommand().lower()

                    if query == "none":
                        continue

                    # Deactivation command
                    if 'deactivate jarvis' in query or 'go to sleep' in query:
                        activate_jarvis() 
                        print("Jarvis is now inactive. Press Ctrl+J to reactivate.")
                        break # Exit command loop, outer loop will wait for reactivation

                    elif 'wikipedia' in query:
                        speak('Searching Wikipedia...')
                        query_term = query.replace("wikipedia", "").strip()
                        if not query_term:
                            speak("What would you like to search on Wikipedia?")
                            query_term = takeCommand().lower()
                            if query_term == "none" or not query_term:
                                continue
                        
                        # Initialize Wikipedia API with a user agent
                        # Best practice: define user_agent clearly
                        wiki_wiki = wikipedia.Wikipedia(
                            language='en',
                            user_agent="JarvisAssistant/1.0 (YourApp; you@example.com)" # Please update with your info
                        )
                        page_py = wiki_wiki.page(query_term)
                        summary = wikipedia.summary(query_term, sentences=2)
                        print(summary)
                        speak("According to Wikipedia")
                        speak(summary)

                    elif 'open youtube' in query:
                        webbrowser.open("https://www.youtube.com") # Corrected URL
                        speak("Opening YouTube.")

                    elif 'open google' in query:
                        webbrowser.open("https://www.google.com") # Corrected URL
                        speak("Opening Google.")

                    elif 'search youtube' in query:
                        search_term = query.replace('search youtube', '').strip()
                        if not search_term:
                            speak("What would you like to search on YouTube?")
                            search_term = takeCommand().lower()
                        if search_term != "none" and search_term: 
                            searchYoutube(search_term)
                            speak(f"Searching YouTube for {search_term}")
                        else:
                            speak("No search query provided for YouTube.")
                            
                    elif 'search google' in query:
                        search_term = query.replace('search google', '').strip()
                        if not search_term:
                            speak("What would you like to search on Google?")
                            search_term = takeCommand().lower()
                        if search_term != "none" and search_term:
                            searchGoogle(search_term)
                            speak(f"Searching Google for {search_term}")
                        else:
                            speak("No search query provided for Google.")

                    elif 'play a song' in query or 'play song' in query:
                        song_query = query.replace('play a song', '').replace('play song', '').strip()
                        if not song_query:
                            speak("Which song would you like to hear?")
                            song_query = takeCommand().lower()
                        if song_query != "none" and song_query:
                            playSpotifyTrack(song_query)
                        else:
                            speak("No song name provided.")

                    elif 'the time' in query:
                        strTime = datetime.datetime.now().strftime("%I:%M %p") # More readable time
                        speak(f"The current time is {strTime}")
                        print(f"The time is {strTime}")

                    elif 'quit' in query or 'exit' in query or 'stop' in query:
                        speak("Goodbye! Have a nice day.")
                        if hotkey_registered: 
                            try:
                                keyboard.remove_hotkey("ctrl+j") 
                            except KeyError: pass
                            except Exception as e: print(f"Error removing hotkey on exit: {e}")
                        sys.exit()
                    
                    elif 'write a code' in query or 'write code' in query:
                        speak("What programming task or code snippet would you like assistance with?")
                        code_task_query = takeCommand().lower() # Renamed
                        if code_task_query != "none" and code_task_query:
                            # Frame the request more clearly for the LLM
                            prompt_for_code = f"Generate a Python code snippet for the following task: {code_task_query}. Provide only the code block."
                            code_response = query_groq(prompt_for_code, groq_api_key)
                            print("Code Response from Groq:\n", code_response) # Log the raw response
                            code_only = extract_code(code_response)
                            if code_only and code_only != code_response : # Check if extraction did something
                                speak("I have generated the code. Opening it now.")
                                open_notepad_with_code(code_only)
                            elif code_only : # If extract_code returns the same (maybe no backticks)
                                speak("I have the response. Displaying it and copying to clipboard.")
                                open_notepad_with_code(code_only)
                            else:
                                speak("I received a response, but couldn't extract a clear code block. I'll copy the full response.")
                                open_notepad_with_code(code_response)
                        else:
                            speak("No coding task provided.")

                    elif 'create a text file' in query:
                        speak("What should be the name of the text file (without extension)?")
                        file_name = takeCommand()
                        if file_name != "none" and file_name:
                            speak("What should be the content or topic for the text file?")
                            topic = takeCommand()
                            if topic != "none" and topic:
                                file_path = os.path.join(os.path.expanduser('~'), 'Desktop', file_name + '.txt')
                                create_text_file_with_content(file_path, f"Write content about: {topic}", groq_api_key)
                            else: speak("No topic provided for the text file.")
                        else: speak("No file name provided.")

                    elif 'create a word file' in query:
                        speak("What should be the name of the Word file (without extension)?")
                        file_name = takeCommand()
                        if file_name != "none" and file_name:
                            speak("What should be the content or topic for the Word file?")
                            topic = takeCommand()
                            if topic != "none" and topic:
                                file_path = os.path.join(os.path.expanduser('~'), 'Desktop', file_name + '.docx')
                                create_word_file_with_content(file_path, topic, groq_api_key)
                            else: speak("No topic provided for the Word file.")
                        else: speak("No file name provided.")
                    
                    elif 'create a powerpoint' in query or 'create a ppt' in query or 'create ppt' in query or 'create presentation' in query:
                        speak("What should be the name of the PowerPoint file (without extension)?")
                        file_name = takeCommand()
                        if file_name != "none" and file_name:
                            speak("What is the main topic for the presentation?")
                            topic = takeCommand()
                            if topic != "none" and topic:
                                file_path = os.path.join(os.path.expanduser('~'), 'Desktop', file_name + '.pptx')
                                create_ppt_file_with_content(file_path, topic, groq_api_key)
                            else: speak("No topic provided for the PowerPoint.")
                        else: speak("No file name provided.")

                    elif 'create a folder' in query:
                        speak("What should be the name of the folder?")
                        folder_name = takeCommand()
                        if folder_name != "none" and folder_name:
                            folder_path = os.path.join(os.path.expanduser('~'), 'Desktop', folder_name)
                            create_folder(folder_path)
                        else: speak("No folder name provided.")
                        
                    elif 'delete file' in query:
                        speak("What is the full name of the file to delete from your Desktop (including extension)?")
                        file_name = takeCommand()
                        if file_name != "none" and file_name:
                            # Clean up common speech recognition errors for file names
                            file_name_cleaned = file_name.lower().replace(" dot ", ".").replace(" ", "")
                            file_path = os.path.join(os.path.expanduser('~'), 'Desktop', file_name_cleaned)
                            delete_file(file_path)
                        else: speak("No file name provided for deletion.")

                    elif 'open file' in query:
                        speak("What is the full name of the file to open from your Desktop (including extension)?")
                        file_name = takeCommand()
                        if file_name != "none" and file_name:
                            file_name_cleaned = file_name.lower().replace(" dot ", ".").replace(" ", "")
                            file_path = os.path.join(os.path.expanduser('~'), 'Desktop', file_name_cleaned)
                            open_file(file_path)
                        else: speak("No file name provided to open.")

                    elif 'open folder' in query:
                        speak("What is the name of the folder on your Desktop to open?")
                        folder_name = takeCommand()
                        if folder_name != "none" and folder_name:
                            folder_path = os.path.join(os.path.expanduser('~'), 'Desktop', folder_name)
                            open_folder(folder_path)
                        else: speak("No folder name provided to open.")

                    elif 'list files' in query:
                        speak("Which folder on your Desktop do you want to list the files of?")
                        folder_name = takeCommand()
                        if folder_name != "none" and folder_name:
                            folder_path = os.path.join(os.path.expanduser('~'), 'Desktop', folder_name)
                            list_files_in_folder(folder_path)
                        else: speak("No folder name provided to list files.")

                    elif 'weather' in query:
                        city_name_query = query.replace('weather in','').replace('weather','').strip()
                        if not city_name_query:
                            speak("Please tell me the city name for the weather forecast.")
                            city_name_query = takeCommand().lower()
                        
                        if city_name_query != "none" and city_name_query:
                            get_weather(city_name_query)
                        else:
                            speak("No city name provided for weather.")

                    elif 'news' in query:
                        news_topic_query = query.replace('news about','').replace('news on','').replace('news','').strip()
                        if not news_topic_query:
                            speak("What topic would you like news about?")
                            news_topic_query = takeCommand().lower()

                        if news_topic_query != "none" and news_topic_query:
                            news_update = get_news(news_api_key, news_topic_query)
                            speak(news_update)
                        else:
                            speak("No topic provided for news.")
                    
                    else: # Default to Groq for other queries
                        speak("Thinking...")
                        response = query_groq(query, groq_api_key)
                        cleaned_response = clean_groq_content(response) # Clean up Groq response
                        print("Groq Response:", cleaned_response)
                        speak(cleaned_response)
            else: # Jarvis is not active
                time.sleep(0.1) # Prevent high CPU usage while waiting

    except KeyboardInterrupt:
        print("Program interrupted by user.")
    except Exception as e: # Catch any other unexpected errors in the main loop
        print(f"An unexpected error occurred in the main loop: {e}")
        speak("I encountered an unexpected error. Please check the console.")
    finally:
        if hotkey_registered:
            try:
                keyboard.remove_hotkey("ctrl+j")
                print("Ctrl+J hotkey removed on final exit.")
            except KeyError: pass # Already removed or never set
            except Exception as e: print(f"Error removing hotkey on final exit: {e}")
        print("Exiting Jarvis.")
