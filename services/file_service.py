from pptx import Presentation
from docx import Document
from io import BytesIO
from config import Config
from services.groq_service import process_groq_query

def generate_ppt(topic):
    # Use your existing PPT generation logic adapted for in-memory creation
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_slide.placeholders[1].text = "Created by JARVIS"
    
    # Get content from Groq
    content = process_groq_query(f"Create 3 bullet points about {topic} for a PowerPoint slide")
    
    # Content slide
    bullet_slide = prs.slides.add_slide(prs.slide_layouts[1])
    bullet_slide.shapes.title.text = topic
    text_frame = bullet_slide.shapes[1].text_frame
    
    for point in content.split('\n')[:3]:
        if point.strip():
            p = text_frame.add_paragraph()
            p.text = point.strip()
            p.level = 0
    
    # Save to in-memory bytes
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    
    return ppt_bytes.getvalue()

def generate_word_file(topic):
    # Use your existing Word generation logic adapted for in-memory creation
    doc = Document()
    
    # Add title
    doc.add_heading(topic, level=0)
    
    # Get content from Groq
    content = process_groq_query(f"Write a 3 paragraph explanation about {topic}")
    doc.add_paragraph(content)
    
    # Save to in-memory bytes
    doc_bytes = BytesIO()
    doc.save(doc_bytes)
    doc_bytes.seek(0)
    
    return doc_bytes.getvalue()

def generate_text_file(topic, content=None):
    if content is None:
        content = process_groq_query(f"Write comprehensive information about {topic}")
    
    # Create in-memory text file
    text_bytes = BytesIO()
    text_bytes.write(content.encode('utf-8'))
    text_bytes.seek(0)
    return text_bytes.getvalue()

def generate_code_file(topic, code=None):
    if code is None:
        code = process_groq_query(f"Write Python code for {topic}. Include detailed comments.")
    
    # Create in-memory code file
    code_bytes = BytesIO()
    code_bytes.write(code.encode('utf-8'))
    code_bytes.seek(0)
    return code_bytes.getvalue()